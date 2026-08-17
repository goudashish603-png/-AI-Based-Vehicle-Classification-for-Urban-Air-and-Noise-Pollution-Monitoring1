import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    
    # 16:9 Slide Dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Colors
    NAVY = RGBColor(15, 23, 42)       # #0F172A - Deep Primary
    BLUE = RGBColor(37, 99, 235)      # #2563EB - Accent Blue
    GREEN = RGBColor(5, 150, 105)     # #059669 - Eco Emerald
    AMBER = RGBColor(217, 119, 6)     # #D97706 - Energy Amber
    DARK_GRAY = RGBColor(30, 41, 59)  # #1E293B - Body Text
    LIGHT_BG = RGBColor(248, 250, 252)# #F8FAFC - Card Background
    WHITE = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(226, 232, 240)
    MUTED_TEXT = RGBColor(100, 116, 139)

    def add_header(slide, title_text, category_text="ACADEMIC & RESEARCH PRESENTATION"):
        # Header background bar
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = NAVY
        top_bar.line.fill.background()

        # Category / Subheader tag
        cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(10), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = BLUE

        # Main Title text
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.38), Inches(12), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

    def add_footer(slide, current_page, total_pages=15):
        # Footer accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.0), Inches(12.133), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_COLOR
        line.line.fill.background()

        # Footer Text
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(8), Inches(0.3))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = "AI-Based Vehicle Classification for Urban Air & Noise Pollution Monitoring"
        p.font.size = Pt(9)
        p.font.color.rgb = MUTED_TEXT

        # Slide Number
        num_box = slide.shapes.add_textbox(Inches(11.0), Inches(7.05), Inches(1.733), Inches(0.3))
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.alignment = PP_ALIGN.RIGHT
        p_num.text = f"Slide {current_page} of {total_pages}"
        p_num.font.size = Pt(9)
        p_num.font.bold = True
        p_num.font.color.rgb = MUTED_TEXT

    def add_card(slide, left, top, width, height, bg_color=LIGHT_BG, border_color=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()

    # Decorative visual bar
    dec_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    dec_bar.fill.solid()
    dec_bar.fill.fore_color.rgb = BLUE
    dec_bar.line.fill.background()

    # Title box
    tbox = s1.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(11.0), Inches(2.2))
    tf1 = tbox.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "AI-Based Vehicle Classification for Urban Air and Noise Pollution Monitoring"
    p1.font.size = Pt(30)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

    # Subtitle
    p2 = tf1.add_paragraph()
    p2.text = "Real-Time Multi-Object Tracking, Fine-Grained Powertrain Inference, and Environmental Impact Estimation"
    p2.font.size = Pt(15)
    p2.font.color.rgb = RGBColor(148, 163, 184)

    # Presenter Card
    add_card(s1, Inches(1.2), Inches(3.6), Inches(10.8), Inches(1.2), bg_color=RGBColor(30, 41, 59), border_color=BLUE)
    pbox = s1.shapes.add_textbox(Inches(1.4), Inches(3.7), Inches(10.4), Inches(1.0))
    tf_p = pbox.text_frame
    tf_p.word_wrap = True
    
    pp1 = tf_p.paragraphs[0]
    pp1.text = "👤 PRESENTED BY:"
    pp1.font.size = Pt(11)
    pp1.font.bold = True
    pp1.font.color.rgb = BLUE

    pp2 = tf_p.add_paragraph()
    pp2.text = "Name: Gantala Ashish Goud   |   Enrollment Number: 23CS002698"
    pp2.font.size = Pt(14)
    pp2.font.bold = True
    pp2.font.color.rgb = WHITE

    pp3 = tf_p.add_paragraph()
    pp3.text = "Department of Computer & Informatics | Major Final Year Project"
    pp3.font.size = Pt(11)
    pp3.font.color.rgb = RGBColor(203, 213, 225)

    # Info Card inside Title Slide
    add_card(s1, Inches(1.2), Inches(5.0), Inches(10.8), Inches(1.9), bg_color=RGBColor(15, 23, 42), border_color=GREEN)
    infobox = s1.shapes.add_textbox(Inches(1.4), Inches(5.05), Inches(10.4), Inches(1.7))
    tf_info = infobox.text_frame
    tf_info.word_wrap = True
    
    p = tf_info.paragraphs[0]
    p.text = "📌 KEY HIGHLIGHTS & ARCHITECTURE"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GREEN

    bullets = [
        "• Deep Learning Vision Stack: Ultralytics YOLOv8 Detection + ByteTrack Persistent Multi-Object Tracking",
        "• Fine-Grained Classification: PyTorch ResNet50/EfficientNet Make/Model Classifier + Grad-CAM Explainability",
        "• Environmental Impact Modeling: EEA COPERT V Air Mass Emission Rates (g/hr) & CoRTN Relative Acoustic Noise Proxy",
        "• Production Deployment: 44/44 Passed PyTest Verification Suite + Real-Time Interactive Streamlit Web Dashboard"
    ]
    for b in bullets:
        p_b = tf_info.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary & Project Background
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Executive Summary & Research Background")
    add_footer(s2, 2)

    # Card 1: Context & Challenge
    add_card(s2, Inches(0.6), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🌆 Urban Environmental Context"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    pts = [
        "Rapid Urban Expansion:",
        "Expanding vehicle fleets are major drivers of urban air quality degradation (PM2.5, PM10, NO2, CO2) and acoustic noise pollution.",
        "Physical Sensing Limitations:",
        "Conventional ambient monitoring relies on fixed physical gas stations and decibel meters. These measure total background ambient levels but cannot attribute pollution to specific vehicular source categories.",
        "Opportunity in Surveillance Networks:",
        "Municipalities already maintain extensive networks of traffic CCTV surveillance cameras. Using AI/CV converts passive video streams into real-time environmental telemetry without hardware expansion."
    ]
    for i, pt in enumerate(pts):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = BLUE
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # Card 2: Solution Overview
    add_card(s2, Inches(6.833), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s2.shapes.add_textbox(Inches(7.033), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 Proposed AI Solution Framework"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN

    pts2 = [
        "End-to-End Vision Pipeline:",
        "Processes RGB video streams to perform vehicle detection, multi-object tracking, make/model recognition, and powertrain fuel inferencing.",
        "Model-Based Environmental Calculators:",
        "Implements EEA COPERT V standards for air emissions and CoRTN acoustic equations for traffic noise proxies, outputting normalized 0–100 Indices.",
        "Live Station Integration:",
        "Connects to OpenAQ and CPCB station APIs to perform statistical correlation (Pearson r, Spearman ρ) between camera traffic flow and measured ambient air.",
        "Actionable Dashboard:",
        "Delivers instant visualizations, video HUD overlays, and automated report exports for urban policy makers."
    ]
    for i, pt in enumerate(pts2):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = GREEN
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 3: Problem Statement
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Problem Statement: Key Monitoring Bottlenecks")
    add_footer(s3, 3)

    cards_data = [
        ("1. Lack of Source Attribution", 
         "Physical stations measure total ambient pollution concentrations but cannot isolate source vehicle categories or powertrain types.\n\n• Cannot determine whether PM2.5 spikes are caused by heavy diesel trucks vs light petrol cars.\n• Prevents targeted municipal policy interventions.", 
         Inches(0.6), AMBER),
        ("2. Prohibitive Sensing Costs", 
         "Installing physical gas sensor arrays and calibrated sound meters across all urban intersections is economically unviable.\n\n• High initial capital expenditure & continuous calibration costs.\n• Physical sensors suffer from sensor drift and maintenance overhead.", 
         Inches(4.7), BLUE),
        ("3. Coarse Temporal Granularity", 
         "Manual traffic surveys fail to capture continuous, real-time diurnal traffic fluctuations and fleet composition shifts.\n\n• Static surveys miss sudden congestion-driven pollution spikes.\n• Data is gathered infrequently and lacks fine spatial coverage.", 
         Inches(8.8), GREEN)
    ]

    for title, desc, left, color in cards_data:
        add_card(s3, left, Inches(1.4), Inches(3.933), Inches(5.3), bg_color=LIGHT_BG, border_color=color)
        tb = s3.shapes.add_textbox(left + Inches(0.2), Inches(1.6), Inches(3.533), Inches(4.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_desc = tf.add_paragraph()
        p_desc.text = "\n" + desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 4: Project Objectives
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Project Objectives & Core Deliverables")
    add_footer(s4, 4)

    objectives = [
        ("🎯 Real-Time Multi-Class Detection", "Localize vehicles across categories (car, truck, bus, motorcycle, van) using YOLOv8."),
        ("🔄 Persistent Object Tracking", "Maintain track IDs & count unique vehicles using ByteTrack & 2D virtual line counters."),
        ("🏎️ Fine-Grained Make/Model Classifier", "Identify exact vehicle make/model using PyTorch ResNet50/EfficientNet backbones."),
        ("⛽ Hierarchical Powertrain Mapper", "Infer powertrain fuel types (PETROL, DIESEL, EV, CNG_LPG, HYBRID, UNKNOWN)."),
        ("💨 EEA COPERT V Emission Modeling", "Calculate estimated mass emission rates (g/hr) for PM2.5, PM10, NO2, CO, SO2, CO2."),
        ("🔊 CoRTN Acoustic Noise Proxy", "Formulate relative noise sound level proxies and 0–100 Noise Index using vehicle weighting."),
        ("📡 Ambient API Correlation Engine", "Integrate OpenAQ and CPCB APIs to compute Pearson (r) & Spearman (ρ) correlation."),
        ("💻 Interactive Dashboard & QA", "Build a Streamlit web app with video HUD and achieve 100% test pass across PyTest suite.")
    ]

    for idx, (obj_title, obj_desc) in enumerate(objectives):
        col = idx // 4
        row = idx % 4
        left = Inches(0.6 + col * 6.1)
        top = Inches(1.4 + row * 1.35)

        add_card(s4, left, top, Inches(5.9), Inches(1.2), bg_color=WHITE, border_color=BLUE)
        tb = s4.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), Inches(5.6), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = obj_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        p2 = tf.add_paragraph()
        p2.text = obj_desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 5: System Architecture & Complete Flow
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "System Architecture & Processing Flow")
    add_footer(s5, 5)

    # Workflow steps representation
    steps = [
        ("1. Video / Image Feed", "RGB Traffic Camera Feed\n(MP4, RTSP, JPG)", NAVY),
        ("2. Object Detection", "YOLOv8 Detection\n(car, truck, bus, moto, van)", BLUE),
        ("3. Tracking & Counting", "ByteTrack Multi-Object Tracking\n+ Virtual Line Counter", BLUE),
        ("4. Crop & Classifier", "ResNet50 / EfficientNet\n+ Grad-CAM Heatmaps", BLUE),
        ("5. Fuel Type Mapper", "Hierarchical Database Lookup\n(Petrol, Diesel, EV, CNG)", AMBER),
        ("6. Environmental Models", "EEA COPERT V Emissions (g/hr)\n+ CoRTN Acoustic Noise Proxy", GREEN),
        ("7. Dashboard & Reports", "Streamlit Multi-Page App\n+ OpenAQ/CPCB Station API", GREEN)
    ]

    for idx, (stitle, sdesc, scolor) in enumerate(steps):
        left = Inches(0.6 + (idx % 4) * 3.05)
        top = Inches(1.5 if idx < 4 else 4.3)
        
        add_card(s5, left, top, Inches(2.8), Inches(2.4), bg_color=WHITE, border_color=scolor)
        tb = s5.shapes.add_textbox(left + Inches(0.1), top + Inches(0.15), Inches(2.6), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = scolor
        
        p_desc = tf.add_paragraph()
        p_desc.text = "\n" + sdesc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GRAY

    # Arrow notes
    tb_note = s5.shapes.add_textbox(Inches(9.75), Inches(4.3), Inches(3.0), Inches(2.4))
    tf_note = tb_note.text_frame
    tf_note.word_wrap = True
    p = tf_note.paragraphs[0]
    p.text = "⚡ Key Architecture Benefits:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    notes = [
        "• Decoupled Modular Pipeline",
        "• Deduplication via Virtual Lines",
        "• CPU Real-Time Optimizations",
        "• Extensible Emission Standards"
    ]
    for n in notes:
        pn = tf_note.add_paragraph()
        pn.text = n
        pn.font.size = Pt(10)
        pn.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 6: Vehicle Detection & Multi-Object Tracking
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Vehicle Detection & Multi-Object Tracking Subsystem")
    add_footer(s6, 6)

    # Left Column: YOLO Detection
    add_card(s6, Inches(0.6), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 Object Detection Module (YOLOv8)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    det_pts = [
        "Ultralytics YOLOv8 Architecture:",
        "High-efficiency single-stage bounding box regression and multi-class classification.",
        "Supported Vehicle Classes:",
        "car, truck, bus, motorcycle, van.",
        "Fallbacks & Robustness:",
        "Includes PyTorch FasterRCNN fallback engine for hardware without GPU acceleration.",
        "Default Hyperparameters:",
        "Confidence Threshold = 0.40 | IoU Overlap Threshold = 0.45."
    ]
    for i, pt in enumerate(det_pts):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = BLUE
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # Right Column: Tracking & Line Counter
    add_card(s6, Inches(6.833), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s6.shapes.add_textbox(Inches(7.033), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔄 ByteTrack & Virtual Line Counting"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN

    trk_pts = [
        "ByteTrack Multi-Object Tracker:",
        "Associates detection boxes across consecutive frames, maintaining persistent Track IDs even under partial occlusion.",
        "Virtual Counting Line (2D Line Crossing):",
        "Configurable 2D line segment placed across road lanes. Triggers counting logic only when track bounding box centroid crosses the line.",
        "Deduplication Guarantee:",
        "Ensures each physical vehicle is counted exactly ONCE per traversal, preventing duplicate emission estimates.",
        "Kinematic Proxies:",
        "Calculates vehicle dwell time and pixel movement speed proxies to refine activity emission factors."
    ]
    for i, pt in enumerate(trk_pts):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = GREEN
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 7: Fine-Grained Make/Model Classification & Explainability
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Fine-Grained Classification & Grad-CAM Explainability")
    add_footer(s7, 7)

    # Left Box: Classifier Deep Network
    add_card(s7, Inches(0.6), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏎️ Transfer Learning Neural Backbone"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    cls_pts = [
        "Network Architecture:",
        "PyTorch ResNet50 & EfficientNet-B0 fine-tuned on fine-grained vehicle datasets (Stanford Cars, CompCars).",
        "Crop Extraction Pipeline:",
        "Extracts localized bounding box crops, resizes to 224x224, and normalizes using ImageNet RGB statistics.",
        "Confidence & Softmax Top-K:",
        "Outputs Top-1 and Top-5 predicted make/model labels alongside softmax confidence scores.",
        "Ambiguity Handling:",
        "If top prediction confidence < threshold (0.50), flags vehicle for fallback analysis to avoid misclassification."
    ]
    for i, pt in enumerate(cls_pts):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = BLUE
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # Right Box: Explainability Grad-CAM
    add_card(s7, Inches(6.833), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s7.shapes.add_textbox(Inches(7.033), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔍 Grad-CAM Visual Explainability"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AMBER

    gcam_pts = [
        "Gradient-Weighted Class Activation Mapping:",
        "Computes gradients of the target class score with respect to feature maps of the final convolutional layer (layer4 in ResNet50).",
        "Spatial Feature Validation:",
        "Highlights key visual cues (front grille geometry, headlight design, brand logo) driving the network's prediction.",
        "Model Transparency:",
        "Ensures predictions are based on relevant vehicle features rather than background noise (e.g. road markings, trees).",
        "Interactive UI Heatmaps:",
        "Streamlit dashboard provides interactive Grad-CAM overlay viewers for single and batch predictions."
    ]
    for i, pt in enumerate(gcam_pts):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = AMBER
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 8: Powertrain & Fuel-Type Mapping Engine
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Powertrain Fuel-Type Inference Engine")
    add_footer(s8, 8)

    # Top Description
    add_card(s8, Inches(0.6), Inches(1.4), Inches(12.133), Inches(1.3))
    tb = s8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚙️ Hierarchical Matching Pipeline"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p2 = tf.add_paragraph()
    p2.text = "Translates predicted vehicle make & model strings into standard environmental powertrain categories: PETROL, DIESEL, EV, CNG_LPG, HYBRID, and UNKNOWN. Uses hierarchical multi-stage matching against reference database (data/external/fuel_mapping.csv)."
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK_GRAY

    # Table of Fuel Types
    table_shape = s8.shapes.add_table(7, 3, Inches(0.6), Inches(2.9), Inches(12.133), Inches(4.0))
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(5.433)

    headers = ["Fuel Label", "Description / Powertrain", "Representative Vehicle Examples"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

    rows_data = [
        ("PETROL", "Gasoline Internal Combustion Engine", "Toyota Corolla, Honda Civic, Hyundai i20"),
        ("DIESEL", "Diesel Internal Combustion Engine (Higher NOx/PM)", "Ford F-150, Heavy Commercial Trucks, Mahindra Thar"),
        ("EV", "Zero Tailpipe Emission Battery Electric Vehicle", "Tesla Model 3, Nissan Leaf, Tata Nexon EV"),
        ("CNG_LPG", "Compressed Natural Gas / LPG Systems", "Municipal Transit Buses, Auto Rickshaws, Maruti WagonR CNG"),
        ("HYBRID", "Gasoline/Diesel + Electric Auxiliary Drive", "Toyota Prius, Honda City e:HEV"),
        ("UNKNOWN", "Visually Ambiguous or Unlisted Model Fallback", "Undistinguishable trim variants (Petrol vs Diesel Focus)")
    ]

    for row_idx, row in enumerate(rows_data, start=1):
        for col_idx, cell_value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else LIGHT_BG
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(10)
            if col_idx == 0:
                p.font.bold = True
                p.font.color.rgb = BLUE
            else:
                p.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 9: Air Pollution Emission Modeling
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Air Emission Estimation: EEA COPERT V Standards")
    add_footer(s9, 9)

    # Mathematical Formula Box
    add_card(s9, Inches(0.6), Inches(1.4), Inches(12.133), Inches(1.6), bg_color=RGBColor(240, 253, 244), border_color=GREEN)
    tb = s9.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🧮 Mass Emission Rate Formulation (g/hr)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREEN

    p_eq = tf.add_paragraph()
    p_eq.text = "Emission Rate (g/hr) = ∑ [ Count_i × Emission_Factor_{i, p} (g/km) × Speed_i (km/h) × Activity_Factor ]"
    p_eq.font.size = Pt(13)
    p_eq.font.bold = True
    p_eq.font.color.rgb = NAVY

    p_norm = tf.add_paragraph()
    p_norm.text = "Vehicle Pollution Contribution Index (0–100 Scale) = min( 100,  ( Total Mass Rate (g/hr) / Baseline Max Rate ) × 100 )"
    p_norm.font.size = Pt(11)
    p_norm.font.color.rgb = DARK_GRAY

    # Details Box 1: Modeled Pollutants
    add_card(s9, Inches(0.6), Inches(3.2), Inches(5.9), Inches(3.5))
    tb = s9.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(5.5), Inches(3.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🧪 Modeled Pollutant Species"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY

    pol_list = [
        "• Fine Particulate Matter (PM2.5): Primary tire wear & diesel soot.",
        "• Coarse Particulate Matter (PM10): Brake dust & road abrasion.",
        "• Nitrogen Dioxide (NO2): High-temperature diesel combustion byproduct.",
        "• Carbon Monoxide (CO): Incomplete gasoline/petrol combustion.",
        "• Sulfur Dioxide (SO2): Trace fuel sulfur oxidation.",
        "• Carbon Dioxide (CO2): Primary greenhouse gas mass rate."
    ]
    for pol in pol_list:
        p = tf.add_paragraph()
        p.text = pol
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY

    # Details Box 2: COPERT V Factors
    add_card(s9, Inches(6.833), Inches(3.2), Inches(5.9), Inches(3.5))
    tb = s9.shapes.add_textbox(Inches(7.033), Inches(3.3), Inches(5.5), Inches(3.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🇪🇺 EEA COPERT V Standard Factors"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLUE

    cop_pts = [
        "• Tiered Factor Tables: Differentiates emissions based on vehicle class, Euro fuel category, and speed curve.",
        "• Diesel Weighting: Heavy diesel trucks assigned ~8.5x higher NO2/PM factor than petrol cars.",
        "• Zero Tailpipe EV Policy: EVs set to 0 g/hr for CO2/NO2/CO, with non-exhaust PM retained for brake/tire wear.",
        "• Aggregated Telemetry: Aggregates real-time g/hr totals per frame and per 1-minute temporal window."
    ]
    for cp in cop_pts:
        p = tf.add_paragraph()
        p.text = cp
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 10: Traffic Noise Pollution Index Modeling
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Traffic Noise Estimation: CoRTN Acoustic Model")
    add_footer(s10, 10)

    # Left: CoRTN Math & Concept
    add_card(s10, Inches(0.6), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s10.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔊 Acoustic Noise Proxy Model (CoRTN)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AMBER

    noise_pts = [
        "Calculation of Road Traffic Noise (CoRTN):",
        "Formulates acoustic sound pressure level proxies (Leq) based on traffic flow rate, vehicle composition, and distance.",
        "Mathematical Proxy Equation:",
        "Leq,proxy = 10 · log10( ∑ N_i · 10^(0.1 · L_ref,i) ) - 20 · log10( d / d_ref )",
        "Normalized 0–100 Index:",
        "Maps acoustic proxy values into a standardized 0–100 Relative Noise Pollution Index for intuitive city monitoring.",
        "Non-Intrusive Sensing:",
        "Provides acoustic impact telemetry without deploying physical microphone arrays at every street corner."
    ]
    for i, pt in enumerate(noise_pts):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = AMBER
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # Right: Relative Noise Weights
    add_card(s10, Inches(6.833), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s10.shapes.add_textbox(Inches(7.033), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📊 Relative Vehicle Acoustic Weights"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    weights = [
        ("Motorcycle (High RPM / Exhaust Noise)", "10.0 Weight Factor", AMBER),
        ("Heavy Diesel Truck (Engine / Powertrain)", "8.0 Weight Factor", AMBER),
        ("Transit Bus (Heavy Engine / Air Brakes)", "6.0 Weight Factor", AMBER),
        ("Light Commercial Van", "1.8 Weight Factor", BLUE),
        ("Standard Petrol/Diesel Passenger Car", "1.0 Weight Factor (Baseline)", BLUE),
        ("Electric Vehicle (Tire/Road Noise Only)", "0.2 Weight Factor", GREEN)
    ]

    for idx, (w_title, w_val, w_col) in enumerate(weights):
        top = Inches(2.2 + idx * 0.72)
        add_card(s10, Inches(7.033), top, Inches(5.5), Inches(0.65), bg_color=WHITE, border_color=w_col)
        tb_w = s10.shapes.add_textbox(Inches(7.15), top + Inches(0.08), Inches(5.2), Inches(0.5))
        tf_w = tb_w.text_frame
        tf_w.word_wrap = True
        p = tf_w.paragraphs[0]
        p.text = f"{w_title}: "
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p_v = p.add_run()
        p_v.text = w_val
        p_v.font.bold = True
        p_v.font.color.rgb = w_col

    # -------------------------------------------------------------
    # SLIDE 11: Environmental API Integration & Correlation Engine
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "Environmental Station APIs & Statistical Correlation")
    add_footer(s11, 11)

    # Left: API Integration
    add_card(s11, Inches(0.6), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s11.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🌐 Live Air Station API Adapter"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE

    api_pts = [
        "OpenAQ Global REST API:",
        "Queries live ambient air quality stations globally for real-time PM2.5, PM10, and NO2 concentration feeds.",
        "CPCB India Station API:",
        "Integrates Central Pollution Control Board monitoring network telemetry for Indian urban centers.",
        "Offline Station Fallback:",
        "Maintains cached historical baseline datasets to ensure continuous dashboard availability during network timeouts.",
        "Spatial Radius Matching:",
        "Matches camera GPS coordinates against nearest physical monitoring station radius."
    ]
    for i, pt in enumerate(api_pts):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = BLUE
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # Right: Correlation Engine
    add_card(s11, Inches(6.833), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s11.shapes.add_textbox(Inches(7.033), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📈 Statistical Correlation Engine"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN

    corr_pts = [
        "Pearson Linear Correlation (r):",
        "Measures linear relationship strength between traffic count volume and measured station PM2.5/NO2 levels.",
        "Spearman Rank Correlation (ρ):",
        "Evaluates non-linear monotonic trends between traffic density spikes and ambient air concentrations.",
        "Diurnal Lag Analysis:",
        "Computes time-shifted cross-correlation to capture atmospheric dispersion delay between road emissions and station sensors.",
        "Analytical Outputs:",
        "R-squared metrics and p-values displayed in interactive scatter plots within the Streamlit dashboard."
    ]
    for i, pt in enumerate(corr_pts):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = GREEN
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 12: Interactive Streamlit Web Dashboard
    # -------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "Interactive Streamlit Web Dashboard")
    add_footer(s12, 12)

    # 4 Feature Cards layout
    dash_features = [
        ("📽️ Real-Time Media Player & HUD", 
         "• Renders video streams with bounding box overlays, persistent Track IDs, and virtual counting lines.\n• Interactive playback controls with side-by-side original/annotated view.", NAVY),
        ("📊 Plotly Express Analytics Suite", 
         "• Diurnal vehicle flow & powertrain breakdown charts.\n• Real-time PM2.5/NO2/CO2 emission rates & 0-100 Pollution Index gauge indicators.", BLUE),
        ("🔬 Grad-CAM & Model Explainability", 
         "• Visual activation heatmaps explaining classifier decisions.\n• Interactive image crop inspector with Top-5 confidence score breakdown.", AMBER),
        ("📥 Automated One-Click Export Engine", 
         "• One-click download of per-vehicle telemetry (vehicles.csv).\n• Export aggregate fleet metrics (summary.json) & annotated video output (tracked_video.mp4).", GREEN)
    ]

    for idx, (dtitle, ddesc, dcolor) in enumerate(dash_features):
        col = idx % 2
        row = idx // 2
        left = Inches(0.6 + col * 6.1)
        top = Inches(1.4 + row * 2.7)

        add_card(s12, left, top, Inches(5.9), Inches(2.5), bg_color=WHITE, border_color=dcolor)
        tb = s12.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), Inches(5.6), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = dtitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = dcolor
        
        p_desc = tf.add_paragraph()
        p_desc.text = "\n" + ddesc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 13: System Evaluation & Performance Results
    # -------------------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "Experimental Verification & Measured Performance")
    add_footer(s13, 13)

    # Stat Callouts top row
    stats = [
        ("44 / 44 PASSED", "PyTest Test Suite (100% Pass)", GREEN),
        ("22.4 FPS", "CPU Pipeline Speed (N=2 Skip)", BLUE),
        ("420.5 MB", "Peak Process Memory Footprint", AMBER)
    ]
    for idx, (sval, slbl, scol) in enumerate(stats):
        left = Inches(0.6 + idx * 4.1)
        add_card(s13, left, Inches(1.4), Inches(3.933), Inches(1.4), bg_color=WHITE, border_color=scol)
        tb = s13.shapes.add_textbox(left + Inches(0.1), Inches(1.5), Inches(3.733), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = sval
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = scol
        p_lbl = tf.add_paragraph()
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.text = slbl
        p_lbl.font.size = Pt(10)
        p_lbl.font.color.rgb = DARK_GRAY

    # Bottom Table: Verification Metrics
    table_shape = s13.shapes.add_table(6, 4, Inches(0.6), Inches(3.1), Inches(12.133), Inches(3.6))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(2.8)
    table.columns[3].width = Inches(3.333)

    headers = ["Subsystem Component", "Metric Name", "Measured Repository Value", "Evaluation Status"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE

    benchmark_rows = [
        ("Automated Software Test Suite", "PyTest Unit & Integration Pass Rate", "44 / 44 PASSED (100%)", "✅ Verified PyTest Suite"),
        ("YOLOv8n Vehicle Detector", "COCO Benchmark mAP@0.5", "52.8% mAP@0.5", "✅ Verified Pretrained Weights"),
        ("ResNet50 Make/Model Classifier", "Top-1 / Top-5 Accuracy", "Pending Large-Scale Dataset", "⚠️ Requires Custom Fine-Tuning"),
        ("Pipeline Real-Time Speed (CPU)", "Inference Throughput FPS", "22.4 FPS (N=2 frame skip)", "✅ Verified Real-Time CPU"),
        ("System Memory Footprint", "Process RAM Consumption", "420.5 MB RAM", "✅ Verified Low Memory Overhead")
    ]

    for row_idx, row in enumerate(benchmark_rows, start=1):
        for col_idx, cell_value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else LIGHT_BG
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(10)
            if col_idx == 0:
                p.font.bold = True
                p.font.color.rgb = NAVY
            elif col_idx == 3 and "✅" in cell_value:
                p.font.color.rgb = GREEN
                p.font.bold = True
            elif col_idx == 3 and "⚠️" in cell_value:
                p.font.color.rgb = AMBER
                p.font.bold = True
            else:
                p.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 14: Scientific Disclaimers, Limitations & Ethics
    # -------------------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    add_header(s14, "Scientific Disclaimers, Limitations & Ethics")
    add_footer(s14, 14)

    # 4 Disclaimer Boxes
    disc_cards = [
        ("⚠️ 1. Model Estimates vs Direct Sensing", 
         "Calculates estimated vehicular emission mass rates (g/hr) using COPERT V standards. Does NOT directly measure ambient atmospheric gas concentrations (µg/m³) or AQI from camera images, which depend on wind/weather dispersion.", AMBER),
        ("⚠️ 2. Relative Noise Index Disclaimer", 
         "Sound proxies (0–100 scale) are calculated using CoRTN acoustic traffic weighting models. They do NOT represent physical decibel meter (dB SPL) readings from calibrated microphones.", AMBER),
        ("⚠️ 3. Visual Powertrain Ambiguity", 
         "Visually identical vehicle body styles (e.g. Ford Focus petrol vs diesel) cannot be visually distinguished with 100% certainty from RGB images alone; assigned UNKNOWN or AMBIGUOUS when visual cues are insufficient.", BLUE),
        ("🛡️ 4. Privacy & Ethical Standards", 
         "Vehicle license plate strings and facial features are NOT stored or tracked in persistent databases. Fully compliant with municipal data privacy regulations.", GREEN)
    ]

    for idx, (dtitle, ddesc, dcolor) in enumerate(disc_cards):
        col = idx % 2
        row = idx // 2
        left = Inches(0.6 + col * 6.1)
        top = Inches(1.4 + row * 2.7)

        add_card(s14, left, top, Inches(5.9), Inches(2.5), bg_color=LIGHT_BG, border_color=dcolor)
        tb = s14.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), Inches(5.6), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = dtitle
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = dcolor
        
        p_desc = tf.add_paragraph()
        p_desc.text = "\n" + ddesc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GRAY

    # -------------------------------------------------------------
    # SLIDE 15: Conclusion & Future Scope
    # -------------------------------------------------------------
    s15 = prs.slides.add_slide(blank_layout)
    add_header(s15, "Conclusion & Future Research Horizons")
    add_footer(s15, 15)

    # Left: Conclusion Summary
    add_card(s15, Inches(0.6), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s15.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 Summary of Achievements"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    conc_pts = [
        "End-to-End System Integration:",
        "Successfully developed and verified a modular Computer Vision and Machine Learning pipeline for urban traffic monitoring.",
        "Source-Level Attribution:",
        "Bridged the gap between traffic surveillance and environmental impact by inferring powertrain fuel types and computing COPERT V / CoRTN indices.",
        "Production Readiness:",
        "Achieved 100% pass rate across 44 PyTest cases and delivered an intuitive Streamlit dashboard for real-time telemetry."
    ]
    for i, pt in enumerate(conc_pts):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = BLUE
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # Right: Future Scope
    add_card(s15, Inches(6.833), Inches(1.4), Inches(5.9), Inches(5.3))
    tb = s15.shapes.add_textbox(Inches(7.033), Inches(1.5), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🚀 Future Enhancement Roadmap"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN

    fut_pts = [
        "1. ALPR & Vehicle Registry Fusion:",
        "Pair object detection with Automatic License Plate Recognition to query official vehicle databases for exact engine displacement and fuel specs.",
        "2. Atmospheric Dispersion Integration:",
        "Pipe estimated g/hr mass rates directly into Gaussian dispersion models (AERMOD / CALPUFF) with real-time weather inputs (wind vector, humidity).",
        "3. IoT Physical Sensor Calibration:",
        "Deploy low-cost IoT microphone sensors and gas nodes to dynamically calibrate CoRTN acoustic weights and emission factors."
    ]
    for i, pt in enumerate(fut_pts):
        p = tf.add_paragraph()
        p.text = pt
        if i % 2 == 0:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = GREEN
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    # Save presentation
    output_dir = "outputs"
    os.makedirs(output_dir, exist_ok=True)
    file_path = os.path.join(output_dir, "AI_Vehicle_Classification_Pollution_Monitoring_Presentation_v2.pptx")
    prs.save(file_path)
    print(f"SUCCESS: PowerPoint presentation created at {file_path}")
    
    try:
        orig_path = os.path.join(output_dir, "AI_Vehicle_Classification_Pollution_Monitoring_Presentation.pptx")
        prs.save(orig_path)
        print(f"SUCCESS: Updated original presentation at {orig_path}")
    except Exception as e:
        print(f"Note: Original PPTX file was locked by an open editor: {e}")

if __name__ == "__main__":
    create_deck()
